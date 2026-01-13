"""
Document chunking for narrative embeddings.

Implements semantic + size limit chunking:
- Split on natural boundaries (paragraphs, sections, slides)
- Enforce max token limit per chunk
- Split oversized sections on sentence boundaries
- Overlap between chunks to preserve context at boundaries
- Section headers as hard boundaries
"""

import re
from dataclasses import dataclass
from pathlib import Path
from typing import Iterator, List, Tuple

# Max chunk size (~1000 tokens, assuming ~4 chars per token)
MAX_CHUNK_CHARS = 4000

# Min chunk size to avoid tiny fragments
MIN_CHUNK_CHARS = 100

# Overlap between chunks (~100 tokens) to preserve context at boundaries
OVERLAP_CHARS = 400

# Section header patterns (hard boundaries - never split across these)
SECTION_HEADER_PATTERNS = [
    r'^#{1,4}\s+.+$',                    # Markdown headers: # Header, ## Header
    r'^[A-Z][A-Z0-9\s\-]{2,50}:?\s*$',   # ALL CAPS HEADERS (with optional colon)
    r'^\d+\.\s+[A-Z].{5,80}$',           # Numbered sections: 1. Introduction
    r'^(?:Section|Part|Chapter)\s+\d+',  # Section 1, Part 2, Chapter 3
    r'^---+\s*$',                        # Horizontal rules
]

# Compiled pattern for efficiency
_SECTION_HEADER_RE = re.compile(
    '|'.join(f'({p})' for p in SECTION_HEADER_PATTERNS),
    re.MULTILINE
)


@dataclass
class Chunk:
    """A chunk of text from a document."""

    text: str
    source_file: str
    chunk_index: int  # 0-indexed sequence within file
    page_number: int  # Original page/section number (for reference)
    total_chunks: int  # Total chunks in this file
    file_type: str

    @property
    def chunk_id(self) -> str:
        """Unique ID for this chunk."""
        safe_name = self.source_file.replace("/", "_").replace("\\", "_")
        return f"{safe_name}__c{self.chunk_index:04d}"


def is_section_header(text: str) -> bool:
    """Check if text is a section header (hard boundary)."""
    text = text.strip()
    if not text or len(text) > 100:  # Headers are short
        return False
    return bool(_SECTION_HEADER_RE.match(text))


# Sentence-ending punctuation pattern
_SENTENCE_END_CHARS = '.!?'
_SENTENCE_END_RE = re.compile(r'[.!?](?:\s|$)')


def ends_at_sentence_boundary(text: str) -> bool:
    """Check if text ends at a sentence boundary."""
    text = text.rstrip()
    return bool(text) and text[-1] in _SENTENCE_END_CHARS


def trim_to_sentence_end(text: str, min_length: int = MIN_CHUNK_CHARS) -> Tuple[str, str]:
    """Trim text back to the last sentence boundary.

    Args:
        text: Text to trim.
        min_length: Minimum length to preserve (won't trim below this).

    Returns:
        Tuple of (trimmed_text, remainder). If no sentence boundary found
        or trimming would go below min_length, returns (text, "").
    """
    text = text.rstrip()

    # Already ends at sentence boundary?
    if ends_at_sentence_boundary(text):
        return text, ""

    # Find the last sentence-ending punctuation
    last_boundary = -1
    for i in range(len(text) - 1, min_length - 1, -1):
        if text[i] in _SENTENCE_END_CHARS:
            # Verify it's a real sentence end (not abbreviation like "Dr.")
            # Simple heuristic: followed by space/newline or end of text
            if i == len(text) - 1 or text[i + 1] in ' \n\t':
                last_boundary = i
                break

    if last_boundary > min_length:
        trimmed = text[:last_boundary + 1]
        remainder = text[last_boundary + 1:].lstrip()
        return trimmed, remainder

    # No good boundary found, return as-is
    return text, ""


def split_into_sentences(text: str) -> List[str]:
    """Split text into sentences."""
    # Split on sentence boundaries (.!? followed by space or newline)
    sentences = re.split(r'(?<=[.!?])\s+', text)
    return [s.strip() for s in sentences if s.strip()]


def get_overlap_text(text: str, max_chars: int = OVERLAP_CHARS) -> str:
    """Extract overlap text from the end of a chunk.

    Ensures the overlap:
    1. Starts at a sentence boundary (complete sentences only)
    2. Ends at a sentence boundary (from the source text)

    Args:
        text: Source text to extract overlap from.
        max_chars: Maximum characters for overlap.

    Returns:
        Clean overlap text starting and ending at sentence boundaries.
    """
    if not text:
        return ""

    text = text.rstrip()

    # If text is short enough, use it all (if it ends cleanly)
    if len(text) <= max_chars:
        if ends_at_sentence_boundary(text):
            return text
        trimmed, _ = trim_to_sentence_end(text)
        return trimmed

    # Take last max_chars as starting point
    overlap = text[-max_chars:]

    # Find start: look for first sentence boundary (after .!?)
    # We want to start at the beginning of a sentence
    sentence_starts = list(re.finditer(r'[.!?]\s+', overlap))
    if sentence_starts:
        # Start after the first sentence-ending punctuation + whitespace
        first_boundary = sentence_starts[0].end()
        if first_boundary < len(overlap) - MIN_CHUNK_CHARS:
            overlap = overlap[first_boundary:]

    # Ensure overlap ends at sentence boundary
    if not ends_at_sentence_boundary(overlap):
        overlap, _ = trim_to_sentence_end(overlap, min_length=MIN_CHUNK_CHARS // 2)

    return overlap.strip()


def split_large_section(text: str, max_chars: int = MAX_CHUNK_CHARS) -> List[str]:
    """Split a large section into smaller chunks on sentence boundaries.

    Ensures each chunk ends at a sentence boundary when possible.
    For extremely long sentences that must be word-split, tries to
    end at clause boundaries (comma, semicolon) as fallback.
    """
    if len(text) <= max_chars:
        # Ensure even short sections end at sentence boundary
        if ends_at_sentence_boundary(text):
            return [text]
        trimmed, remainder = trim_to_sentence_end(text)
        if remainder:
            return [trimmed, remainder] if len(remainder) >= MIN_CHUNK_CHARS else [text]
        return [text]

    sentences = split_into_sentences(text)
    chunks = []
    current_chunk = []
    current_length = 0

    for sentence in sentences:
        sentence_length = len(sentence)

        # If a single sentence is too long, split by words at clause boundaries
        if sentence_length > max_chars:
            if current_chunk:
                chunks.append(" ".join(current_chunk))
                current_chunk = []
                current_length = 0

            # Split long sentence, preferring clause boundaries
            word_chunks = _split_long_sentence(sentence, max_chars)
            # Add all but last as complete chunks
            chunks.extend(word_chunks[:-1])
            # Last part continues as current chunk
            if word_chunks:
                last_part = word_chunks[-1]
                current_chunk = [last_part]
                current_length = len(last_part)
            continue

        # Check if adding this sentence would exceed limit
        if current_length + sentence_length + 1 > max_chars and current_chunk:
            chunks.append(" ".join(current_chunk))
            current_chunk = []
            current_length = 0

        current_chunk.append(sentence)
        current_length += sentence_length + 1

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    return [c for c in chunks if len(c) >= MIN_CHUNK_CHARS]


def _split_long_sentence(sentence: str, max_chars: int) -> List[str]:
    """Split an extremely long sentence, preferring clause boundaries.

    Tries to split at clause boundaries (;,) before falling back to
    word boundaries. This preserves more readable chunks.
    """
    # First try splitting on semicolons (strong clause boundary)
    if ';' in sentence:
        clauses = sentence.split(';')
        if all(len(c.strip()) <= max_chars for c in clauses):
            # Rejoin with semicolons, keeping punctuation
            return [c.strip() + (';' if i < len(clauses) - 1 else '')
                    for i, c in enumerate(clauses) if c.strip()]

    # Fall back to word-based splitting with comma preference
    words = sentence.split()
    chunks = []
    current_words = []
    current_length = 0

    for word in words:
        word_len = len(word) + 1  # +1 for space

        if current_length + word_len > max_chars and current_words:
            chunk_text = " ".join(current_words)

            # Try to end at a comma if we're mid-sentence
            if not chunk_text.rstrip().endswith((',', ';', ':')):
                # Look for last comma in chunk
                last_comma = chunk_text.rfind(',')
                if last_comma > len(chunk_text) // 2:  # Only if comma is in latter half
                    # Split at comma, carry remainder forward
                    chunks.append(chunk_text[:last_comma + 1])
                    remainder_words = chunk_text[last_comma + 1:].split()
                    current_words = remainder_words + [word]
                    current_length = sum(len(w) + 1 for w in current_words)
                    continue

            chunks.append(chunk_text)
            current_words = []
            current_length = 0

        current_words.append(word)
        current_length += word_len

    if current_words:
        chunks.append(" ".join(current_words))

    return chunks


def extract_paragraphs(text: str) -> List[str]:
    """Extract paragraphs from text, splitting on double newlines."""
    # Split on double newlines or slide markers
    paragraphs = re.split(r'\n\s*\n|(?=^## Slide \d+)', text, flags=re.MULTILINE)
    return [p.strip() for p in paragraphs if p.strip()]


def semantic_chunk(
    paragraphs: List[str],
    max_chars: int = MAX_CHUNK_CHARS,
    overlap_chars: int = OVERLAP_CHARS
) -> List[str]:
    """Create semantic chunks from paragraphs with size limits and overlap.

    Args:
        paragraphs: List of paragraphs to chunk.
        max_chars: Maximum characters per chunk.
        overlap_chars: Characters to overlap between chunks.

    Returns:
        List of chunk texts with overlap applied.
    """
    if not paragraphs:
        return []

    # First pass: create raw chunks respecting section boundaries
    raw_chunks = []
    current_chunk = []
    current_length = 0

    for para in paragraphs:
        para_length = len(para)
        para_is_header = is_section_header(para)

        # Section header forces a new chunk (hard boundary)
        if para_is_header and current_chunk:
            raw_chunks.append("\n\n".join(current_chunk))
            current_chunk = []
            current_length = 0

        # If paragraph is too large, split it
        if para_length > max_chars:
            # First, flush current chunk
            if current_chunk:
                raw_chunks.append("\n\n".join(current_chunk))
                current_chunk = []
                current_length = 0

            # Split the large paragraph
            sub_chunks = split_large_section(para, max_chars)
            raw_chunks.extend(sub_chunks)
            continue

        # Check if adding this paragraph would exceed limit
        if current_length + para_length + 2 > max_chars and current_chunk:
            raw_chunks.append("\n\n".join(current_chunk))
            current_chunk = []
            current_length = 0

        current_chunk.append(para)
        current_length += para_length + 2  # +2 for \n\n

    if current_chunk:
        raw_chunks.append("\n\n".join(current_chunk))

    # Filter out tiny chunks
    raw_chunks = [c for c in raw_chunks if len(c) >= MIN_CHUNK_CHARS]

    if not raw_chunks:
        return []

    # Second pass: add overlap from previous chunk
    # Skip overlap if chunk starts with a section header (clean break)
    final_chunks = [raw_chunks[0]]  # First chunk has no previous overlap

    for i in range(1, len(raw_chunks)):
        current = raw_chunks[i]
        previous = raw_chunks[i - 1]

        # Don't add overlap if current chunk starts with a section header
        # (indicates intentional clean break)
        first_line = current.split('\n')[0].strip()
        if is_section_header(first_line):
            final_chunks.append(current)
            continue

        # Get overlap from previous chunk
        overlap = get_overlap_text(previous, overlap_chars)

        if overlap and len(overlap) + len(current) <= max_chars * 1.2:
            # Add overlap with a marker for debugging (can be removed)
            final_chunks.append(f"{overlap}\n\n{current}")
        else:
            final_chunks.append(current)

    return final_chunks


def chunk_pdf(filepath: Path) -> Iterator[Chunk]:
    """Extract and semantically chunk PDF content."""
    import fitz  # PyMuPDF

    doc = fitz.open(str(filepath))

    # Extract all text with page markers
    all_text_parts = []
    page_boundaries = []  # Track where each page starts in terms of paragraphs

    for page_num, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()
        if text:
            # Store page boundary info
            page_boundaries.append({
                "page": page_num,
                "start_para": len(all_text_parts)
            })
            # Extract paragraphs from this page
            paragraphs = extract_paragraphs(text)
            all_text_parts.extend(paragraphs)

    doc.close()

    if not all_text_parts:
        return

    # Create semantic chunks
    chunks_text = semantic_chunk(all_text_parts)

    # Yield chunks with metadata
    total_chunks = len(chunks_text)
    for chunk_index, text in enumerate(chunks_text):
        # Estimate which page this chunk came from (rough approximation)
        # This is imperfect but gives a useful reference
        page_num = 1
        for boundary in page_boundaries:
            if chunk_index >= boundary["start_para"]:
                page_num = boundary["page"]

        yield Chunk(
            text=text,
            source_file=filepath.name,
            chunk_index=chunk_index,
            page_number=page_num,
            total_chunks=total_chunks,
            file_type="pdf",
        )


def chunk_docx(filepath: Path) -> Iterator[Chunk]:
    """Extract and semantically chunk DOCX content."""
    from docx import Document

    doc = Document(str(filepath))

    # Collect all text blocks
    paragraphs = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)

    for table in doc.tables:
        table_text = []
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip(" |"):
                table_text.append(row_text)
        if table_text:
            paragraphs.append("\n".join(table_text))

    if not paragraphs:
        return

    # Create semantic chunks
    chunks_text = semantic_chunk(paragraphs)

    total_chunks = len(chunks_text)
    for chunk_index, text in enumerate(chunks_text):
        yield Chunk(
            text=text,
            source_file=filepath.name,
            chunk_index=chunk_index,
            page_number=chunk_index + 1,  # Section number
            total_chunks=total_chunks,
            file_type="docx",
        )


def chunk_txt(filepath: Path) -> Iterator[Chunk]:
    """Chunk text files semantically."""
    text = filepath.read_text(encoding="utf-8", errors="replace")

    paragraphs = extract_paragraphs(text)

    if not paragraphs:
        return

    chunks_text = semantic_chunk(paragraphs)

    total_chunks = len(chunks_text)
    for chunk_index, chunk_text in enumerate(chunks_text):
        yield Chunk(
            text=chunk_text,
            source_file=filepath.name,
            chunk_index=chunk_index,
            page_number=chunk_index + 1,
            total_chunks=total_chunks,
            file_type="txt",
        )


def chunk_xlsx(filepath: Path) -> Iterator[Chunk]:
    """Extract and chunk XLSX content by sheet."""
    import pandas as pd

    xlsx = pd.ExcelFile(filepath)

    all_chunks = []
    for sheet_name in xlsx.sheet_names:
        try:
            df = pd.read_excel(xlsx, sheet_name=sheet_name)
            if not df.empty:
                text = f"Sheet: {sheet_name}\n\n{df.to_string(index=False)}"

                # Split large sheets
                if len(text) > MAX_CHUNK_CHARS:
                    sub_chunks = split_large_section(text)
                    for sub in sub_chunks:
                        all_chunks.append((sheet_name, sub))
                else:
                    all_chunks.append((sheet_name, text))
        except Exception:
            continue

    total_chunks = len(all_chunks)
    for chunk_index, (sheet_name, text) in enumerate(all_chunks):
        yield Chunk(
            text=text,
            source_file=filepath.name,
            chunk_index=chunk_index,
            page_number=chunk_index + 1,  # Sheet number
            total_chunks=total_chunks,
            file_type="xlsx",
        )


def chunk_pptx(filepath: Path) -> Iterator[Chunk]:
    """Extract and chunk PPTX content by slide."""
    from pptx import Presentation

    prs = Presentation(str(filepath))

    all_chunks = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    para_text = "".join(run.text for run in para.runs)
                    if para_text.strip():
                        texts.append(para_text.strip())

            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        texts.append(row_text)

        if texts:
            slide_text = "\n".join(texts)

            # Split large slides
            if len(slide_text) > MAX_CHUNK_CHARS:
                sub_chunks = split_large_section(slide_text)
                for sub in sub_chunks:
                    all_chunks.append((slide_num, sub))
            else:
                all_chunks.append((slide_num, slide_text))

    total_chunks = len(all_chunks)
    for chunk_index, (slide_num, text) in enumerate(all_chunks):
        yield Chunk(
            text=text,
            source_file=filepath.name,
            chunk_index=chunk_index,
            page_number=slide_num,
            total_chunks=total_chunks,
            file_type="pptx",
        )


def chunk_document(filepath: Path) -> Iterator[Chunk]:
    """Chunk a document based on its file type."""
    ext = filepath.suffix.lower()

    try:
        if ext == ".pdf":
            yield from chunk_pdf(filepath)
        elif ext in {".docx", ".doc"}:
            yield from chunk_docx(filepath)
        elif ext == ".txt":
            yield from chunk_txt(filepath)
        elif ext in {".xlsx", ".xls"}:
            yield from chunk_xlsx(filepath)
        elif ext == ".pptx":
            yield from chunk_pptx(filepath)
        else:
            # Skip unsupported formats
            pass
    except Exception as e:
        print(f"  Error chunking {filepath.name}: {e}")
