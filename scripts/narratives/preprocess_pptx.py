#!/usr/bin/env python3
"""
Extract text content from PPTX files for narrative processing.

Converts PPTX files to text files that can be processed by the narratives pipeline.
"""

import argparse
import json
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_text_from_shape(shape) -> str:
    """Extract text from a shape, handling different shape types."""
    text_parts = []

    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            para_text = ""
            for run in paragraph.runs:
                para_text += run.text
            if para_text.strip():
                text_parts.append(para_text.strip())

    # Handle tables
    if shape.has_table:
        table = shape.table
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    row_text.append(cell_text)
            if row_text:
                text_parts.append(" | ".join(row_text))

    # Handle grouped shapes
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            sub_text = extract_text_from_shape(sub_shape)
            if sub_text:
                text_parts.append(sub_text)

    return "\n".join(text_parts)


def extract_pptx_text(pptx_path: Path) -> dict:
    """Extract all text from a PPTX file."""
    prs = Presentation(str(pptx_path))

    result = {
        "source_file": pptx_path.name,
        "slide_count": len(prs.slides),
        "slides": []
    }

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data = {
            "slide_number": slide_num,
            "content": []
        }

        for shape in slide.shapes:
            text = extract_text_from_shape(shape)
            if text:
                slide_data["content"].append(text)

        if slide_data["content"]:
            result["slides"].append(slide_data)

    return result


def format_as_markdown(data: dict) -> str:
    """Format extracted data as markdown for processing."""
    lines = [
        f"# {data['source_file']}",
        f"",
        f"**Slides:** {data['slide_count']}",
        f"",
    ]

    for slide in data["slides"]:
        lines.append(f"## Slide {slide['slide_number']}")
        lines.append("")
        for content in slide["content"]:
            lines.append(content)
            lines.append("")

    return "\n".join(lines)


def process_directory(input_dir: Path, output_dir: Path, dry_run: bool = False) -> dict:
    """Process all PPTX files in a directory."""
    pptx_files = list(input_dir.glob("*.pptx"))

    stats = {
        "total": len(pptx_files),
        "processed": 0,
        "skipped": 0,
        "errors": []
    }

    output_dir.mkdir(parents=True, exist_ok=True)

    for pptx_path in pptx_files:
        output_path = output_dir / f"{pptx_path.stem}.txt"

        if output_path.exists():
            print(f"  Skipping (exists): {pptx_path.name}")
            stats["skipped"] += 1
            continue

        if dry_run:
            print(f"  Would process: {pptx_path.name}")
            stats["processed"] += 1
            continue

        try:
            print(f"  Processing: {pptx_path.name}")
            data = extract_pptx_text(pptx_path)
            markdown = format_as_markdown(data)

            output_path.write_text(markdown, encoding="utf-8")
            stats["processed"] += 1

        except Exception as e:
            print(f"  ERROR: {pptx_path.name}: {e}")
            stats["errors"].append({"file": pptx_path.name, "error": str(e)})

    return stats


def main():
    parser = argparse.ArgumentParser(description="Extract text from PPTX files")
    parser.add_argument("input_dir", type=Path, nargs="?", help="Directory containing PPTX files")
    parser.add_argument("--output-dir", type=Path, help="Output directory (default: input_dir/extracted)")
    parser.add_argument("--dry-run", action="store_true", help="Show what would be processed")
    parser.add_argument("--single", type=Path, help="Process single file")

    args = parser.parse_args()

    if args.single:
        # Single file mode
        data = extract_pptx_text(args.single)
        print(format_as_markdown(data))
        return

    output_dir = args.output_dir or args.input_dir / "extracted"

    print(f"Input: {args.input_dir}")
    print(f"Output: {output_dir}")
    print()

    stats = process_directory(args.input_dir, output_dir, args.dry_run)

    print()
    print(f"Summary: {stats['processed']} processed, {stats['skipped']} skipped, {len(stats['errors'])} errors")

    if stats["errors"]:
        print("\nErrors:")
        for err in stats["errors"]:
            print(f"  {err['file']}: {err['error']}")


if __name__ == "__main__":
    main()
