"""
MXI PowerPoint Presentation Helpers
====================================
Reusable functions for creating MXI-branded presentations
for the Samsung Taylor FAB1 delay analysis project.

Usage:
    from pptx_helpers import MXIPresentation

    prs = MXIPresentation("Analysis Report", "Scope Growth Findings")
    prs.add_content_slide("Key Findings", ["Finding 1", "Finding 2"])
    prs.add_chart_slide("Trend Analysis", fig)
    prs.save("output.pptx")
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import List, Optional, Tuple, Union
import os


# =============================================================================
# MXI Asset Paths (OneDrive)
# =============================================================================
class MXIAssets:
    """Paths to MXI brand assets."""
    BASE_DIR = Path("/mnt/c/Users/pcuri/OneDrive - MXI/Desktop/Samsung Dashboard/UI")

    # Logos
    LOGO_TRANSPARENT = BASE_DIR / "MXI Logo - Transparent.png"
    LOGO_FULL = BASE_DIR / "MXI Logo Full.png"  # Use only on white backgrounds
    LOGO = BASE_DIR / "MXI Logo.png"
    SAMSUNG_LOGO = BASE_DIR / "Samsung Logo.png"

    # Banners and Backgrounds
    BANNER = BASE_DIR / "MXI Banner PNG.png"  # Gradient banner for headers
    LIGHT_BACKGROUND = BASE_DIR / "MXI Light Background.png"

    @classmethod
    def verify_assets(cls) -> dict:
        """Check which assets are available."""
        assets = {
            "logo_transparent": cls.LOGO_TRANSPARENT.exists(),
            "logo_full": cls.LOGO_FULL.exists(),
            "logo": cls.LOGO.exists(),
            "samsung_logo": cls.SAMSUNG_LOGO.exists(),
            "banner": cls.BANNER.exists(),
            "light_background": cls.LIGHT_BACKGROUND.exists(),
        }
        return assets


# MXI Brand Colors
class MXIColors:
    NAVY = RGBColor(0, 51, 102)        # Primary - headers, title bars
    BLUE = RGBColor(0, 112, 192)       # Secondary accent
    ACCENT = RGBColor(48, 75, 106)     # #304B6A - Subtle accent for lines
    GRAY = RGBColor(89, 89, 89)        # Body text
    DARK_GRAY = RGBColor(64, 64, 64)   # Darker text for better contrast
    LIGHT_GRAY = RGBColor(128, 128, 128)  # Footers
    WHITE = RGBColor(255, 255, 255)    # Title text on dark backgrounds
    SAMSUNG_BLUE = RGBColor(20, 40, 160)  # Samsung branding
    RED = RGBColor(192, 0, 0)          # Alerts, negative trends
    GREEN = RGBColor(0, 128, 0)        # Success, positive trends
    YELLOW = RGBColor(255, 192, 0)     # Warnings, caution


class MXIPresentation:
    """Create MXI-branded PowerPoint presentations."""

    def __init__(self, title: str, subtitle: str = "", date: str = None,
                 show_logo: bool = True, show_samsung_logo: bool = False):
        """
        Initialize a new MXI presentation.

        Args:
            title: Main presentation title
            subtitle: Optional subtitle
            date: Date string (defaults to current month/year)
            show_logo: Include MXI logo on slides (default True)
            show_samsung_logo: Include Samsung logo for co-branding (default False)
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # Widescreen 16:9
        self.prs.slide_height = Inches(7.5)

        self.title = title
        self.subtitle = subtitle
        self.date = date or datetime.now().strftime("%B %Y")
        self.slide_count = 0

        # Logo settings
        self.show_logo = show_logo
        self.show_samsung_logo = show_samsung_logo

        # Verify assets exist
        if self.show_logo and not MXIAssets.LOGO_FULL.exists():
            print(f"Warning: MXI logo not found at {MXIAssets.LOGO_FULL}")
            self.show_logo = False

        # Add title slide
        self._add_title_slide()

    def _add_title_slide(self):
        """Add the title slide with MXI branding, banner, and logo."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Banner image as header (gradient) - taller for title slide
        header_height = Inches(2.8)
        if MXIAssets.BANNER.exists():
            slide.shapes.add_picture(
                str(MXIAssets.BANNER),
                Inches(0), Inches(0),
                width=self.prs.slide_width,
                height=header_height
            )
        else:
            # Fallback to solid navy
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0,
                self.prs.slide_width, header_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = MXIColors.NAVY
            shape.line.fill.background()
            shape.shadow.inherit = False

        # MXI Logo (centered at top of banner)
        if self.show_logo and MXIAssets.LOGO_TRANSPARENT.exists():
            # Center the logo: (slide_width - logo_width) / 2
            # Logo height 0.6", estimate width ~1.8" based on aspect ratio
            logo_width_estimate = Inches(1.8)
            logo_left = (self.prs.slide_width - logo_width_estimate) / 2
            slide.shapes.add_picture(
                str(MXIAssets.LOGO_TRANSPARENT),
                logo_left, Inches(0.25),
                height=Inches(0.6)
            )

        # Customer logo (left of title, if enabled)
        title_left = Inches(0.75)
        if self.show_samsung_logo and MXIAssets.SAMSUNG_LOGO.exists():
            slide.shapes.add_picture(
                str(MXIAssets.SAMSUNG_LOGO),
                Inches(0.5), Inches(1.15),
                height=Inches(0.55)
            )
            title_left = Inches(2.2)  # Shift title right to accommodate customer logo

        # Title
        title_box = slide.shapes.add_textbox(
            title_left, Inches(1.1), Inches(10.5), Inches(0.9)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = self.title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = MXIColors.WHITE
        p.font.name = "Segoe UI"

        # Subtitle
        if self.subtitle:
            sub_box = slide.shapes.add_textbox(
                title_left, Inches(1.85), Inches(10.5), Inches(0.6)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = self.subtitle
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(180, 195, 210)
            p.font.name = "Segoe UI"

        # Project info box (below banner, on white background)
        info_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(3.8), Inches(6), Inches(1.5)
        )
        tf = info_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Samsung Taylor FAB1"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = MXIColors.NAVY
        p.font.name = "Segoe UI"

        p = tf.add_paragraph()
        p.text = "Performance Analysis"
        p.font.size = Pt(14)
        p.font.color.rgb = MXIColors.GRAY
        p.font.name = "Segoe UI"

        # Date and prepared by (no bottom MXI logo - minimalistic)
        footer_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(6.8), Inches(8), Inches(0.4)
        )
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{self.date}  |  Prepared by MXI"
        p.font.size = Pt(12)
        p.font.color.rgb = MXIColors.LIGHT_GRAY
        p.font.name = "Segoe UI"

        self.slide_count += 1
        return slide

    def _add_slide_title(self, slide, title: str):
        """Add a minimalist content title - no header banner, just clean typography."""
        # Title text - navy color on white background
        title_box = slide.shapes.add_textbox(
            Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.7)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = MXIColors.NAVY
        p.font.name = "Segoe UI"

        # Accent line under title - spans most of width
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.9), Inches(1.1),
            Inches(11.5), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = MXIColors.ACCENT
        line.line.fill.background()
        line.shadow.inherit = False

    def _add_slide_footer(self, slide):
        """Add banner footer with logo to a content slide."""
        self.slide_count += 1

        # Banner at bottom
        banner_height = Inches(0.65)
        banner_top = self.prs.slide_height - banner_height

        if MXIAssets.BANNER.exists():
            slide.shapes.add_picture(
                str(MXIAssets.BANNER),
                Inches(0), banner_top,
                width=self.prs.slide_width,
                height=banner_height
            )
        else:
            # Fallback to solid color
            banner = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), banner_top,
                self.prs.slide_width, banner_height
            )
            banner.fill.solid()
            banner.fill.fore_color.rgb = MXIColors.NAVY
            banner.line.fill.background()
            banner.shadow.inherit = False

        # MXI Logo on right side of banner - far right
        if self.show_logo and MXIAssets.LOGO_TRANSPARENT.exists():
            slide.shapes.add_picture(
                str(MXIAssets.LOGO_TRANSPARENT),
                Inches(12.3), banner_top + Inches(0.1),
                height=Inches(0.45)
            )

        # Page number at beginning of footer (semibold)
        page_box = slide.shapes.add_textbox(
            Inches(0.5), banner_top + Inches(0.18), Inches(0.5), Inches(0.35)
        )
        tf = page_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(self.slide_count)
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(200, 210, 220)
        p.font.name = "Segoe UI Semibold"

        # Footer text on banner (semibold) - includes project, report title, and date
        footer = slide.shapes.add_textbox(
            Inches(1.0), banner_top + Inches(0.18), Inches(10), Inches(0.35)
        )
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.text = f"MXI  |  Samsung Taylor FAB1  |  {self.title}  |  {self.date}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(200, 210, 220)
        p.font.name = "Segoe UI Semibold"

    def add_content_slide(self, title: str, bullets: List[str],
                          sub_bullets: dict = None) -> 'Slide':
        """
        Add a content slide with bullet points.

        Args:
            title: Slide title
            bullets: List of bullet point strings
            sub_bullets: Dict mapping bullet index to list of sub-bullets
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_title(slide, title)

        # Content area - starts below title with generous white space
        content_box = slide.shapes.add_textbox(
            Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.2)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        sub_bullets = sub_bullets or {}

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Add bullet character for better visibility
            p.text = f"•  {bullet}"
            p.font.size = Pt(15)
            p.font.color.rgb = MXIColors.DARK_GRAY
            p.font.name = "Segoe UI"
            p.space_before = Pt(5)
            p.space_after = Pt(5)
            p.level = 0

            # Add sub-bullets if present
            if i in sub_bullets:
                for sub in sub_bullets[i]:
                    sp = tf.add_paragraph()
                    sp.text = f"    ‒  {sub}"  # Indented dash for sub-bullets
                    sp.font.size = Pt(13)
                    sp.font.color.rgb = MXIColors.GRAY
                    sp.font.name = "Segoe UI"
                    sp.level = 1
                    sp.space_before = Pt(3)
                    sp.space_after = Pt(3)

        self._add_slide_footer(slide)
        return slide

    def add_two_column_slide(self, title: str,
                             left_title: str, left_content: List[str],
                             right_title: str, right_content: List[str]) -> 'Slide':
        """Add a two-column content slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_title(slide, title)

        # Column dimensions - generous margins and center gap
        left_margin = Inches(0.9)
        col_width = Inches(5.5)
        center_gap = Inches(0.9)
        right_start = left_margin + col_width + center_gap

        # Left column title
        left_title_box = slide.shapes.add_textbox(
            left_margin, Inches(1.6), col_width, Inches(0.5)
        )
        tf = left_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = MXIColors.BLUE
        p.font.name = "Segoe UI"

        # Left column content
        left_box = slide.shapes.add_textbox(
            left_margin, Inches(2.2), col_width, Inches(4.5)
        )
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(left_content):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = MXIColors.GRAY
            p.font.name = "Segoe UI"
            p.space_before = Pt(1)
            p.space_after = Pt(1)

        # Right column title
        right_title_box = slide.shapes.add_textbox(
            right_start, Inches(1.6), col_width, Inches(0.5)
        )
        tf = right_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = MXIColors.BLUE
        p.font.name = "Segoe UI"

        # Right column content
        right_box = slide.shapes.add_textbox(
            right_start, Inches(2.2), col_width, Inches(4.5)
        )
        tf = right_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(right_content):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = MXIColors.GRAY
            p.font.name = "Segoe UI"
            p.space_before = Pt(1)
            p.space_after = Pt(1)

        self._add_slide_footer(slide)
        return slide

    def add_chart_slide(self, title: str, fig,
                        caption: str = None,
                        left: float = 0.9, top: float = 1.5,
                        width: float = 11.5, height: float = 5.0) -> 'Slide':
        """
        Add a slide with a matplotlib figure.

        Args:
            title: Slide title
            fig: matplotlib Figure object
            caption: Optional caption below the chart
            left, top, width, height: Position in inches
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_title(slide, title)

        # Save figure to buffer
        image_stream = BytesIO()
        fig.savefig(image_stream, format='png', dpi=150,
                    bbox_inches='tight', facecolor='white', edgecolor='none')
        image_stream.seek(0)

        # Add image
        slide.shapes.add_picture(
            image_stream,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )

        # Add caption if provided
        if caption:
            caption_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.3)
            )
            tf = caption_box.text_frame
            p = tf.paragraphs[0]
            p.text = caption
            p.font.size = Pt(10)
            p.font.italic = True
            p.font.color.rgb = MXIColors.LIGHT_GRAY
            p.font.name = "Segoe UI"
            p.alignment = PP_ALIGN.CENTER

        self._add_slide_footer(slide)
        return slide

    def add_table_slide(self, title: str, data: List[List],
                        col_widths: List[float] = None) -> 'Slide':
        """
        Add a slide with a table.

        Args:
            title: Slide title
            data: 2D list where first row is headers
            col_widths: Optional list of column widths in inches
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_title(slide, title)

        rows = len(data)
        cols = len(data[0])

        # Calculate table dimensions - generous margins
        table_width = Inches(11.5)
        table_height = Inches(min(5.0, rows * 0.45))

        table = slide.shapes.add_table(
            rows, cols,
            Inches(0.9), Inches(1.5),
            table_width, table_height
        ).table

        # Set column widths if provided
        if col_widths:
            for i, width in enumerate(col_widths):
                table.columns[i].width = Inches(width)

        # Style header row - minimalistic light design
        for j in range(cols):
            cell = table.cell(0, j)
            cell.text = str(data[0][j])
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light gray header

            p = cell.text_frame.paragraphs[0]
            p.font.color.rgb = MXIColors.DARK_GRAY
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.name = "Segoe UI"
            p.alignment = PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Style data rows - subtle alternating
        for i in range(1, rows):
            for j in range(cols):
                cell = table.cell(i, j)
                cell.text = str(data[i][j])

                # Very subtle alternating row colors
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(252, 252, 252)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(11)
                p.font.color.rgb = MXIColors.DARK_GRAY
                p.font.name = "Segoe UI"
                p.alignment = PP_ALIGN.LEFT
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        self._add_slide_footer(slide)
        return slide

    def add_kpi_slide(self, title: str, kpis: List[dict]) -> 'Slide':
        """
        Add a KPI dashboard slide.

        Args:
            title: Slide title
            kpis: List of dicts with keys: label, value, trend (up/down/flat), color (optional)

        Example:
            kpis = [
                {"label": "Schedule Slip", "value": "127 days", "trend": "down"},
                {"label": "Scope Growth", "value": "+34%", "trend": "down"},
                {"label": "Issues Resolved", "value": "847", "trend": "up"},
            ]
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_title(slide, title)

        # Calculate positions for KPI boxes - generous spacing
        num_kpis = len(kpis)
        box_width = min(2.6, 11 / num_kpis - 0.4)
        gap = 0.5  # Gap between boxes
        total_width = num_kpis * box_width + (num_kpis - 1) * gap
        start_x = (13.333 - total_width) / 2

        for i, kpi in enumerate(kpis):
            x = start_x + i * (box_width + gap)

            # KPI box background - flat, no shadow
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(1.8),
                Inches(box_width), Inches(4.0)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(248, 249, 250)
            box.line.color.rgb = RGBColor(230, 230, 230)
            box.shadow.inherit = False  # Disable shadow

            # KPI label
            label_box = slide.shapes.add_textbox(
                Inches(x + 0.15), Inches(2.2),
                Inches(box_width - 0.3), Inches(0.6)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = kpi["label"]
            p.font.size = Pt(13)
            p.font.color.rgb = MXIColors.GRAY
            p.font.name = "Segoe UI"
            p.alignment = PP_ALIGN.CENTER

            # KPI value
            value_box = slide.shapes.add_textbox(
                Inches(x + 0.15), Inches(3.2),
                Inches(box_width - 0.3), Inches(1.0)
            )
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(kpi["value"])
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.name = "Segoe UI"
            p.alignment = PP_ALIGN.CENTER

            # Determine color based on trend
            trend = kpi.get("trend", "flat")
            if "color" in kpi:
                p.font.color.rgb = kpi["color"]
            elif trend == "up":
                p.font.color.rgb = MXIColors.GREEN
            elif trend == "down":
                p.font.color.rgb = MXIColors.RED
            else:
                p.font.color.rgb = MXIColors.NAVY

            # Trend indicator
            trend_box = slide.shapes.add_textbox(
                Inches(x + 0.15), Inches(4.6),
                Inches(box_width - 0.3), Inches(0.5)
            )
            tf = trend_box.text_frame
            p = tf.paragraphs[0]
            if trend == "up":
                p.text = "▲"
                p.font.color.rgb = MXIColors.GREEN
            elif trend == "down":
                p.text = "▼"
                p.font.color.rgb = MXIColors.RED
            else:
                p.text = "●"
                p.font.color.rgb = MXIColors.GRAY
            p.font.size = Pt(20)
            p.alignment = PP_ALIGN.CENTER

        self._add_slide_footer(slide)
        return slide

    def add_section_slide(self, section_title: str, section_number: int = None) -> 'Slide':
        """Add a section divider slide - title with number, accent line, and banner at bottom."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Section number and title on same line
        title_top = Inches(2.8)
        if section_number:
            # Number in accent color
            num_text = f"0{section_number}" if section_number < 10 else str(section_number)
            full_title = f"{num_text}   {section_title}"
        else:
            full_title = section_title

        # Section title - navy on white
        title_box = slide.shapes.add_textbox(
            Inches(0.9), title_top, Inches(11.5), Inches(1.2)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = full_title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = MXIColors.NAVY
        p.font.name = "Segoe UI"

        # Accent line under title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.9), title_top + Inches(1.0),
            Inches(5.0), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = MXIColors.ACCENT
        line.line.fill.background()
        line.shadow.inherit = False

        # Banner at bottom
        banner_height = Inches(0.65)
        banner_top = self.prs.slide_height - banner_height

        if MXIAssets.BANNER.exists():
            slide.shapes.add_picture(
                str(MXIAssets.BANNER),
                Inches(0), banner_top,
                width=self.prs.slide_width,
                height=banner_height
            )
        else:
            banner = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), banner_top,
                self.prs.slide_width, banner_height
            )
            banner.fill.solid()
            banner.fill.fore_color.rgb = MXIColors.NAVY
            banner.line.fill.background()
            banner.shadow.inherit = False

        # MXI Logo on right side of banner
        if self.show_logo and MXIAssets.LOGO_TRANSPARENT.exists():
            slide.shapes.add_picture(
                str(MXIAssets.LOGO_TRANSPARENT),
                Inches(11.5), banner_top + Inches(0.1),
                height=Inches(0.45)
            )

        self.slide_count += 1
        return slide

    def save(self, filename: str = None, output_dir: str = None):
        """
        Save the presentation.

        Args:
            filename: Output filename (auto-generated if not provided)
            output_dir: Output directory (defaults to OneDrive presentations folder)
        """
        if output_dir is None:
            output_dir = "/mnt/c/Users/pcuri/OneDrive - MXI/Desktop/Samsung Dashboard/Presentations"

        # Create directory if it doesn't exist
        Path(output_dir).mkdir(parents=True, exist_ok=True)

        if filename is None:
            date_str = datetime.now().strftime("%Y%m%d")
            safe_title = self.title.lower().replace(" ", "_")[:30]
            filename = f"{date_str}_{safe_title}.pptx"

        if not filename.endswith(".pptx"):
            filename += ".pptx"

        filepath = os.path.join(output_dir, filename)
        self.prs.save(filepath)
        print(f"Presentation saved: {filepath}")
        return filepath


# Convenience function for quick presentations
def create_quick_presentation(title: str, slides_data: List[dict],
                               output_path: str = None) -> str:
    """
    Create a presentation quickly from structured data.

    Args:
        title: Presentation title
        slides_data: List of dicts with 'type' and content keys
        output_path: Optional output path

    Example:
        slides_data = [
            {"type": "content", "title": "Overview", "bullets": ["Point 1", "Point 2"]},
            {"type": "section", "title": "Analysis", "number": 1},
            {"type": "kpis", "title": "Metrics", "kpis": [...]},
        ]
    """
    prs = MXIPresentation(title)

    for slide_data in slides_data:
        slide_type = slide_data.get("type", "content")

        if slide_type == "content":
            prs.add_content_slide(
                slide_data["title"],
                slide_data.get("bullets", []),
                slide_data.get("sub_bullets")
            )
        elif slide_type == "section":
            prs.add_section_slide(
                slide_data["title"],
                slide_data.get("number")
            )
        elif slide_type == "kpis":
            prs.add_kpi_slide(
                slide_data["title"],
                slide_data.get("kpis", [])
            )
        elif slide_type == "table":
            prs.add_table_slide(
                slide_data["title"],
                slide_data.get("data", [[]]),
                slide_data.get("col_widths")
            )
        elif slide_type == "two_column":
            prs.add_two_column_slide(
                slide_data["title"],
                slide_data.get("left_title", ""),
                slide_data.get("left_content", []),
                slide_data.get("right_title", ""),
                slide_data.get("right_content", [])
            )

    return prs.save(output_path)


if __name__ == "__main__":
    # Example usage
    prs = MXIPresentation(
        "Samsung Taylor FAB1 Analysis",
        "Q4 2025 Progress Report"
    )

    prs.add_content_slide("Executive Summary", [
        "Schedule analysis complete for 66 YATES snapshots",
        "Identified 127-day slip from baseline schedule",
        "Scope grew by 2,847 tasks (34% increase)",
        "1,108 issues documented and categorized"
    ])

    prs.add_section_slide("Key Findings", 1)

    prs.add_kpi_slide("Performance Metrics", [
        {"label": "Schedule Slip", "value": "127 days", "trend": "down"},
        {"label": "Scope Growth", "value": "+34%", "trend": "down"},
        {"label": "Issues Found", "value": "1,108", "trend": "flat"},
        {"label": "Data Sources", "value": "66", "trend": "up"}
    ])

    prs.add_table_slide("Delay Attribution", [
        ["Category", "Days", "% of Total", "Primary Cause"],
        ["Coordination", "45", "35%", "Trade conflicts"],
        ["Rework", "32", "25%", "Quality issues"],
        ["Design Changes", "28", "22%", "Client requests"],
        ["Weather", "12", "9%", "Rain delays"],
        ["Other", "10", "8%", "Various"]
    ])

    prs.add_two_column_slide(
        "Next Steps",
        "Completed", [
            "Schedule variance analysis",
            "Issue categorization",
            "Labor correlation study",
            "Weekly report extraction"
        ],
        "In Progress", [
            "Quality impact assessment",
            "Cost reconciliation",
            "Final report preparation",
            "Client presentation"
        ]
    )

    # Save to current directory for testing
    prs.save("example_presentation.pptx", output_dir=".")
